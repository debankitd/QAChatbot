import streamlit as st
import faiss
import os
from io import BytesIO
from docx import Document
import numpy as np
from PyPDF2 import PdfReader
from pptx import Presentation
from dotenv import load_dotenv
from langchain.chains import RetrievalQA, LLMChain
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.embeddings.openai import OpenAIEmbeddings
from langchain.vectorstores import FAISS
from langchain.docstore.in_memory import InMemoryDocstore
from langchain.chat_models import ChatOpenAI
from langchain.prompts import PromptTemplate
import pdfplumber
from PIL import Image
import pytesseract

# ============================================================
#  Environment Setup
# ============================================================
load_dotenv()
openai_api_key = os.getenv("OPENAI_API_KEY")

if not openai_api_key:
    st.error("⚠️ OpenAI API Key not found. Add it to your .env file.")
    st.stop()

USE_SUMMARIZER = True
SUMMARY_THRESHOLD = 1500

# ============================================================
#  Utility Functions
# ============================================================

def summarize_text_if_needed(text, threshold=SUMMARY_THRESHOLD):
    """Summarize text if it exceeds the given threshold."""
    if not USE_SUMMARIZER or len(text) <= threshold:
        return text

    summarizer_llm = ChatOpenAI(model="gpt-3.5-turbo", temperature=0.3, openai_api_key=openai_api_key)
    prompt = PromptTemplate(
        input_variables=["content"],
        template=("Summarize the following text clearly and concisely, keeping all key facts:\n\n{content}\n\nSummary:")
    )
    chain = LLMChain(llm=summarizer_llm, prompt=prompt)
    summary = chain.run(content=text)
    return summary.strip()


def extract_text_from_pdf(file_bytes):
    """Extract text from a PDF; fallback to OCR if necessary."""
    text = ""
    try:
        with pdfplumber.open(BytesIO(file_bytes)) as pdf:
            for page in pdf.pages:
                page_text = page.extract_text()
                if page_text:
                    text += page_text + "\n"
        if not text.strip():  # fallback if no text layer found
            st.warning("⚙️ No extractable text found, using OCR...")
            text = extract_text_from_image_pdf(file_bytes)
    except Exception as e:
        st.warning(f"PDF parsing error: {e}. Trying OCR fallback...")
        text = extract_text_from_image_pdf(file_bytes)
    return text.strip()


def extract_text_from_image_pdf(file_bytes):
    """Extract text from scanned PDF using OCR (pytesseract)."""
    text = ""
    try:
        with pdfplumber.open(BytesIO(file_bytes)) as pdf:
            for page in pdf.pages:
                img = page.to_image(resolution=200).original
                ocr_text = pytesseract.image_to_string(Image.fromarray(img))
                text += ocr_text + "\n"
    except Exception as e:
        st.error(f"OCR failed: {e}")
    return text.strip()


def extract_text_from_file(uploaded_file):
    """Extracts text content from multiple file formats."""
    ext = uploaded_file.name.split(".")[-1].upper()
    file_bytes = uploaded_file.read()  # Read once
    text = ""

    try:
        if ext == "PDF":
            text = extract_text_from_pdf(file_bytes)
        elif ext == "DOCX":
            doc = Document(BytesIO(file_bytes))
            text = " ".join([p.text for p in doc.paragraphs if p.text])
        elif ext == "TXT":
            text = file_bytes.decode("utf-8")
        elif ext == "PPTX":
            ppt = Presentation(BytesIO(file_bytes))
            slides_text = []
            for slide in ppt.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        slides_text.append(shape.text)
            text = " ".join(slides_text)
        else:
            st.warning(f"⚠️ Unsupported file type: {uploaded_file.name}")
    except Exception as e:
        st.error(f"Error reading {uploaded_file.name}: {e}")

    return text.strip()


def build_vectorstore_from_files(files):
    """Build FAISS vectorstore from uploaded documents."""
    all_text = []
    for f in files:
        extracted = extract_text_from_file(f)
        if extracted:
            all_text.append(extracted)

    if not all_text:
        raise ValueError("No text extracted from files.")

    st.info(f"🧠 Creating vectorstore from {len(all_text)} file(s)...")

    text_splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=200)
    chunks = text_splitter.split_text(" ".join(all_text))

    embeddings = OpenAIEmbeddings(openai_api_key=openai_api_key)
    dim = len(embeddings.embed_query("sample text"))
    index = faiss.IndexFlatL2(dim)

    vs = FAISS(
        embedding_function=embeddings.embed_query,
        index=index,
        docstore=InMemoryDocstore(),
        index_to_docstore_id={}
    )
    for i in range(0, len(chunks), 50):
        vs.add_texts(chunks[i:i+50])

    st.success(f"✅ Indexed {len(chunks)} chunks successfully.")
    return vs


def build_context_from_history(messages, max_chars=3000):
    """Summarize or trim chat history for context injection."""
    full_history = "\n".join([f"{m['role'].capitalize()}: {m['content']}" for m in messages])
    if len(full_history) > max_chars:
        return summarize_text_if_needed(full_history)
    return full_history

def answer_question(vectorstore, query, context):
    """Ask GPT-4 with retrieval and context injection."""
    llm = ChatOpenAI(model="gpt-4", openai_api_key=openai_api_key, temperature=0.2)
    retriever = vectorstore.as_retriever(search_kwargs={"k": 2})

    system_prompt = (
        "You are an intelligent data assistant. Use the provided context and retrieved document snippets "
        "to answer the user question accurately.\n\n"
        f"Conversation context:\n{context}\n\n"
        "Relevant retrieved information:\n{context}\n\n"
        "User question: {question}\nAssistant:"
    )

    qa_prompt = PromptTemplate(
        template=system_prompt,
        input_variables=["context", "question"]
    )

    chain = RetrievalQA.from_chain_type(
        llm=llm,
        retriever=retriever,
        chain_type="stuff",
        chain_type_kwargs={"prompt": qa_prompt},
    )

    return chain.run({"query": query})


# ============================================================
#  Streamlit App
# ============================================================
def main():
    st.set_page_config(page_title="📚 AI Data Assistant", layout="wide")

    st.markdown("<h1 style='text-align:center;'>🤖 AI Data Assistant</h1>", unsafe_allow_html=True)
    st.caption("Chat with your uploaded documents — memory powered by Streamlit session state.")

    if "messages" not in st.session_state:
        st.session_state["messages"] = []

    with st.sidebar:
        st.header("📤 Upload Your Documents")
        uploaded_files = st.file_uploader("Upload PDFs, Word, Text, or PPTX", type=["pdf", "docx", "txt", "pptx"], accept_multiple_files=True)
        if st.button("⚡ Process Files", use_container_width=True):
            if not uploaded_files:
                st.warning("Please upload at least one document.")
            else:
                try:
                    vectorstore = build_vectorstore_from_files(uploaded_files)
                    st.session_state["vectorstore"] = vectorstore
                    st.success("✅ Files processed successfully!")
                    st.session_state["messages"].append({"role": "assistant", "content": "Documents indexed. You can start asking questions now."})
                except Exception as e:
                    st.error(f"Error: {e}")

    # Chat UI
    if "vectorstore" in st.session_state:
        for msg in st.session_state["messages"]:
            avatar = "🤖" if msg["role"] == "assistant" else "👤"
            st.markdown(f"<div style='margin-bottom:0.6rem;'>{avatar} {msg['content']}</div>", unsafe_allow_html=True)

        if user_input := st.chat_input("Ask your question here..."):
            st.session_state["messages"].append({"role": "user", "content": user_input})
            st.chat_message("user").write(user_input)

            try:
                context = build_context_from_history(st.session_state["messages"])
                response = answer_question(st.session_state["vectorstore"], user_input, context)
                st.session_state["messages"].append({"role": "assistant", "content": response})
                st.chat_message("assistant").write(response)
            except Exception as e:
                st.error(f"Error generating answer: {e}")


if __name__ == "__main__":
    main()
